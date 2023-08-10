import pandas as pd
import numpy as np
import streamlit as st
from kiwipiepy import Kiwi
from pptx import Presentation

st.title("🖨 발표용 프롬프트 만들기")
st.info('###### 언제 사용하나요?\n강의 준비할 때, 긴 문단을 입력하면 적당히 문장별로 끊어서 PPT로 만들어 드립니다.')
st.warning('###### 어떻게 해결하나요?\n텍스트 입력 ➡ 문장별로 끊어서 슬라이드로!')

output_file_name = st.text_input("출력할 파일명을 입력해주세요. ", 'prompt')
max_number = st.number_input("한 슬라이드에 몇자 이내로 들어가게 할까요?", value = 100)

# 1. 문단을 문장으로 자르기 (Kiwi)
def txt_2_sentence_Kiwi(para):
  kiwi = Kiwi()
  split_list = []
  for s in kiwi.split_into_sents(para):
    split_list.append(s.text)
  return split_list

# 2. 100 초과하는 문장은 쉼표 기준으로 자르기
def less_than_100(sentence_list):
  less_than_100 = []
  for s in sentence_list : 
    if len(s)>max_number:
      for s_comma in s.split(','):
        less_than_100.append(s_comma)
    else:
      less_than_100.append(s)
  return less_than_100

# 3. 너무 길이가 짧은 문장은 두개씩 합치기 코드
def almost_100(sentence_list):
  sentence_list_revised = []
  i=0
  while i < len(sentence_list)-1:
    if len(sentence_list[i])+len(sentence_list[i+1])>max_number:
      sentence_list_revised.append(sentence_list[i])
      i = i + 1
    else:
      sentence_list_revised.append(sentence_list[i]+" "+sentence_list[i+1])
      i = i + 2
  sentence_list_revised.append(sentence_list[-1]) #마지막 문장 # 이렇게 하니까 마지막 두 문장의 글자수 합이 100이하일 때 마지막 말이 두번 반복됨. WHY? ######################
  return sentence_list_revised

prompt = st.text_area("Prompt 입력창")
if prompt:
    st.session_state['prompt'] = prompt



if st.button("문장 슬라이스 미리보기"):
    # real code
    sents = txt_2_sentence_Kiwi(prompt)
    sents = less_than_100(sents)
    for i in range(5):
        sents = almost_100(sents)
    for i in range(len(sents)):
       st.write("### slide",i)
       st.write(sents[i])




if st.button("이대로 PPT로 변환하기"):
    prs = Presentation()
    # real code
    sents = txt_2_sentence_Kiwi(prompt)
    sents = less_than_100(sents)
    for i in range(5):
        sents = almost_100(sents)


    # 한 유형의 슬라이드 
    for i in range(len(sents)):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        body = slide.placeholders[1]
        frame = body.text_frame
        frame.text = sents[i]


    # save the presentation to a BytesIO object
    import io
    ppt_bytes = io.BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)

    st.download_button(label="완성! 눌러서 파일 다운로드 받기", data=ppt_bytes, file_name=output_file_name+".pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")