import streamlit as st
import pptx

# 페이지 설명 부분
st.title("🖨PPT reader")


col1, col2 = st.columns(2)
with col1:
    st.info('###### 언제 사용하나요?\nPPT파일을 열어서 참고하기 귀찮을 때, 정리되지 않은 PPT 파일을 한번에 출력하고 싶을 때! 모든 텍스트를 읽어 한꺼번에 출력해드립니다.')
with col2:
    st.warning('###### 어떻게 해결하나요?\nPPT 파일 업로드 ➡ 슬라이드별 내용 출력 ')

st.write("※ 업로드하신 PPT는 따로 저장되지 않으므로 걱정하지 않으셔도 됩니다. ")

def parse_pptx(pptx_file):
    prs = pptx.Presentation(pptx_file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                st.write(text_frame.text.replace('.','.\n'))
                                    
def parse_ppt_sep(pptx_file):
    prs = pptx.Presentation(pptx_file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame

                if text_frame.text == "":
                    continue

                if shape.is_placeholder:
                    # 제목 슬라이드의 경우
                    st.markdown("### " + text_frame.text)
                else:
                    # 본문 슬라이드의 경우
                    paragraphs = text_frame.text.split("\n")
                    formatted_text = "\n".join(paragraphs)
                    formatted_text = formatted_text.replace(".", ".\n")
                    st.write(formatted_text)

# Streamlit 앱 설정
#st.set_page_config(page_title="PPTX Viewer")

# 업로드된 파일 가져오기
st.write("### pptx 파일 업로드")

uploaded_file = st.file_uploader("파일을 먼저 업로드하세요.", type=["pptx"])


if uploaded_file is not None:
    st.write('* 모든 텍스트상자를 구분없이 출력하려면 **구분없이 출력하기**를 클릭해주세요')
    st.write('* 제목과 본문 텍스트상자를 구분하려면 **제목 본문 구분하여 출력하기**를 클릭해주세요')

    # "파싱 시작" 버튼 추가
    if st.checkbox("구분없이 출력하기"):
        # PPTX 파일 파싱하여 출력
        parse_pptx(uploaded_file)
    elif st.checkbox("제목 본문 구분하여 출력하기"):
        parse_ppt_sep(uploaded_file)
