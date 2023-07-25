import pandas as pd
import numpy as np
import streamlit as st
from kiwipiepy import Kiwi
from pptx import Presentation

output_file_name = st.text_input("출력할 파일명을 입력해주세요. ", 'prompt')
#prompt = st.text_input("프롬프트를 입력해주세요. ")
prompt = '안뇽? 저는 제목에서 보셨던 것처럼 반복적인 일을 좋아하지 않는,저는 제목에서 보셨던 것처럼 반복적인 일을 좋아하지 않는,저는 제목에서 보셨던 것처럼 반복적인 일을 좋아하지 않는, 효율적인 것을 좋아하는 사람입니다. 저는 2017년부터 교직 생활을 시작했고, 현재는 서울에 있는 한 중학교에서 수학을 가르치고 있는 교사입니다. 학부때는 수학을 전공하고 수학교육학 석사를 작년에 졸업했으며, 현재는 잠시 휴직을 하고 같은 전공 박사과정 중에 있습니다. 인공지능 기반의 수학교육 평가와 데이터 분석에 관심이 많습니다. 프로그래밍 언어는 학부때 잠깐 씨언어를 배운 것 외에는 아무 것도 모르는, 비전공자입니다. 이런 제가 2020년도부터 파이썬을 처음 접하면서 삶이 많이 달라지게 되었는데요, 코딩을 즐겨하는 선생님들과 함께 만든 커뮤니티 “쪼랩”의 운영진이기도 합니다.'
# 1. 문단을 문장으로 자르기 (Kiwi)
def txt_2_sentence_Kiwi(para):
  kiwi = Kiwi()
  split_list = []
  for s in kiwi.split_into_sents(para):
    split_list.append(s.text)
  return split_list

# 2. 100 초과하는 문장은 쉼표 기준으로 자르기
def less_than_100(sentence_list):
  less_than_100 = []
  for s in sentence_list : 
    if len(s)>100:
      for s_comma in s.split(','):
        less_than_100.append(s_comma)
    else:
      less_than_100.append(s)
  return less_than_100

# 3. 너무 길이가 짧은 문장은 두개씩 합치기 코드
def almost_100(sentence_list):
  sentence_list_revised = []
  i=0
  while i < len(sentence_list)-1:
    if len(sentence_list[i])+len(sentence_list[i+1])>100:
      sentence_list_revised.append(sentence_list[i])
      i = i + 1
    else:
      sentence_list_revised.append(sentence_list[i]+" "+sentence_list[i+1])
      i = i + 2
  sentence_list_revised.append(sentence_list[-1]) #마지막 문장 # 이렇게 하니까 마지막 두 문장의 글자수 합이 100이하일 때 마지막 말이 두번 반복됨. WHY? ######################
  return sentence_list_revised

prompt = st.text_area("Prompt 입력창")
if prompt:
    st.session_state['prompt'] = prompt

if st.button("PPT로 변환하기"):
    prs = Presentation()
    # real code
    sents = txt_2_sentence_Kiwi(prompt)
    sents = less_than_100(sents)
    for i in range(5):
        sents = almost_100(sents)


    # 한 유형의 슬라이드 
    for i in range(len(sents)):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        body = slide.placeholders[1]
        frame = body.text_frame
        frame.text = sents[i]


    # save the presentation to a BytesIO object
    import io
    ppt_bytes = io.BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)

    st.download_button(label="눌러서 파일 다운로드 받기", data=ppt_bytes, file_name=output_file_name+".pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")